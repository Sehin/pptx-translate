from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.slide import Slide


class PPTXConverter:
    def __init__(self, path):
        """
        Путь к pptx файлу
        :param path:
        """
        self.path = path
        self.presentation = Presentation(path)

    def translate(self, dictionary):
        slide_count = len(self.presentation.slides)
        counter = 1
        for slide in self.presentation.slides:
            print(f"SLIDE {counter} of {slide_count}")
            for shape in slide.shapes:
                if shape.has_text_frame:    # Если есть текст в форме
                    self._translate_text_frame(shape, dictionary)
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            self._translate_text_frame(cell, dictionary)
                        pass
                    print("THIS IS TABLE")
            counter += 1

    def _translate_text_frame(self, shape, dictionary):
        # Прохожусь по словарю, если есть совпадения в текст фрейме
        # for key in dictionary.keys() or key.lower() in dictionary.keys():
        for key in dictionary.keys():
            if key in shape.text_frame.text or key.lower() in shape.text_frame.text:
                text_frame = shape.text_frame
                cur_text = text_frame.text
                if len(cur_text) < len(shape.text):
                    new_text = str(dictionary[key])
                else:
                    if key in cur_text:
                        new_text = cur_text.replace(str(key).strip(), str(dictionary[key]))
                    else:
                        new_text = cur_text.replace(str(key.lower()).strip(), str(dictionary[key].lower()))
                print(f"Change {text_frame.text} ||| {new_text}")

                # find not empty run
                for paragraph in list(text_frame.paragraphs):
                    if len(paragraph.runs) > 0:
                        break

                # text_frame.paragraphs[0].runs[0].text = new_text
                paragraph.runs[0].text = new_text.strip()

                if len(shape.text_frame.paragraphs) > 1 or len(shape.text_frame.paragraphs[0].runs) > 0:
                    is_first = True
                    if not paragraph._element is list(text_frame.paragraphs)[0]._element:
                        is_first = False
                    for paragraph_ in shape.text_frame.paragraphs:
                        for run in paragraph_.runs:
                            if is_first or run._r is paragraph.runs[0]._r:
                                is_first = False
                                continue
                            run.text = ''
                if hasattr(shape, "width"):
                    new_text_len = len(paragraph.runs[0].text)
                    old_text_len = len(cur_text)
                    if len(paragraph.runs[0].text) > len(cur_text):
                        delta = new_text_len/old_text_len
                        shape.width = int(shape.width * (delta*1.5))
                        pass

    def save(self, path):
        self.presentation.save(path)

    def search_and_replace(self, search_str, repl_str, input, output):
        """"search and replace text in PowerPoint while preserving formatting"""
        # Useful Links ;)
        # https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
        # https://stackoverflow.com/questions/45247042/how-to-keep-original-text-formatting-of-text-with-python-powerpoint
        prs = self.presentation
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if (shape.text.find(search_str)) != -1:
                        text_frame = shape.text_frame
                        cur_text = text_frame.paragraphs[0].runs[0].text
                        new_text = cur_text.replace(str(search_str), str(repl_str))
                        text_frame.paragraphs[0].runs[0].text = new_text
        # prs.save(output)